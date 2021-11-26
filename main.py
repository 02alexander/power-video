#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches
from glob import glob

files = glob("imgs/*")
keys = [s[s.find('/')+1:s.find('.')] for s in files]
#print(keys)
srtd = sorted(files, key=lambda s: int(s[s.find('/')+1:s.find('.')]))
#print(srtd)

prs = Presentation()

lyt = prs.slide_layouts[6]
for path in srtd:
    slide = prs.slides.add_slide(lyt)
    left=Inches(0)
    top=Inches(0)
    height = Inches(7)
    width = Inches(9)
    img=slide.shapes.add_picture(path,left,top, height=height, width=width)

prs.save("pres.pptx")
